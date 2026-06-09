#!/usr/bin/env python3
"""Generate SnapChef milestone-report PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

# ---- palette --------------------------------------------------------------
BG      = RGBColor(0x0D,0x0D,0x0D)
CARD    = RGBColor(0x1A,0x1A,0x1A)
CARD2   = RGBColor(0x22,0x22,0x22)
TEAL    = RGBColor(0x00,0xC8,0x96)
ORANGE  = RGBColor(0xFF,0x6B,0x35)
TEXT    = RGBColor(0xF0,0xF0,0xF0)
SUB     = RGBColor(0x9A,0xA0,0xA6)
GREEN_F = RGBColor(0x16,0x34,0x1F); GREEN_B = RGBColor(0x2F,0xAE,0x6B); GREEN=RGBColor(0x3F,0xC9,0x7E)
YEL_F   = RGBColor(0x3A,0x31,0x16); YEL_B   = RGBColor(0xE0,0xB3,0x3A); YELLOW=RGBColor(0xE6,0xC25A>>0,0x00) if False else RGBColor(0xE6,0xC2,0x5A)
RED_F   = RGBColor(0x3A,0x17,0x17); RED_B   = RGBColor(0xFF,0x4D,0x4D); RED=RGBColor(0xFF,0x6B,0x6B)
GREY    = RGBColor(0x88,0x88,0x88)
PHBG    = RGBColor(0x3A,0x31,0x16)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return box, tf

def run(p, text, size, color=TEXT, bold=False, italic=False, font="Calibri"):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return r

def para(tf, first=False):
    return tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()

def line1(s, x, y, w, text, size, color, bold=False, align=PP_ALIGN.LEFT, h=0.5):
    box, tf = tb(s, x, y, w, h)
    p = tf.paragraphs[0]; p.alignment = align
    run(p, text, size, color, bold)
    return box

def kicker_title(s, kick, title, title_color=TEAL):
    line1(s, 0.6, 0.35, 12, kick.upper(), 13, ORANGE, bold=True)
    box, tf = tb(s, 0.6, 0.62, 12.1, 1.0)
    p = tf.paragraphs[0]
    run(p, title, 30, title_color, bold=True)
    return box

def card(s, x, y, w, h, fill=CARD, border=RGBColor(0x2A,0x2A,0x2A), border_w=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(border_w)
    sh.shadow.inherit = False
    # soften rounded corner radius
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    return sh

def set_dash(shape):
    ln = shape.line._get_or_add_ln()
    for d in ln.findall(parse_xml('<a:prstDash %s/>'%nsdecls('a')).tag):
        ln.remove(d)
    ln.append(parse_xml('<a:prstDash %s val="dash"/>'%nsdecls('a')))

def add_arrow(shape):
    ln = shape.line._get_or_add_ln()
    ln.append(parse_xml('<a:tailEnd %s type="triangle" w="med" len="med"/>'%nsdecls('a')))

def connector(s, x1,y1,x2,y2, color=GREY, width=1.75, dash=False, arrow=True):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    c.shadow.inherit = False
    if dash: set_dash(c)
    if arrow: add_arrow(c)
    return c

def block(s, x, y, w, h, title, sub, status):
    fills = {"g":(GREEN_F,GREEN_B), "y":(YEL_F,YEL_B), "r":(RED_F,RED_B)}
    f,b = fills[status]
    sh = card(s, x, y, w, h, fill=f, border=b, border_w=1.75)
    tf = sh.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(0.08); tf.margin_right=Inches(0.08)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    run(p, title, 13.5, TEXT, bold=True)
    if sub:
        p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
        run(p2, sub, 10.5, RGBColor(0xC8,0xC8,0xC8))
    return sh

def ph(p, text):
    r = run(p, " "+text+" ", 13, RGBColor(0xFF,0xDF,0x8A), bold=True)
    rPr = r._r.get_or_add_rPr()
    rPr.append(parse_xml('<a:highlight %s><a:srgbClr val="3A3116"/></a:highlight>'%nsdecls('a')))
    return r

def footer(s, n):
    line1(s, 0.6, 7.06, 4, "SnapChef · GIX 2026 Spring", 10, RGBColor(0x66,0x66,0x66))
    line1(s, 12.0, 7.06, 1.0, str(n), 10, RGBColor(0x66,0x66,0x66), align=PP_ALIGN.RIGHT)

def bullets(s, x, y, w, h, items, size=14.5, color=TEXT, gap=6):
    box, tf = tb(s, x, y, w, h)
    for i,(txt, bold) in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.level=0
        run(p, "•  ", size, ORANGE, bold=True)
        # allow simple **bold** lead
        run(p, txt, size, color, bold=bold)
    return box

def style_table(tbl, header=True, font=11):
    # remove default banding/style by setting cells manually
    for ri,row in enumerate(tbl.rows):
        for ci,cell in enumerate(row.cells):
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if (ri==0 and header) else (CARD2 if ri%2 else RGBColor(0x18,0x18,0x18))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left=Inches(0.08); cell.margin_right=Inches(0.08)
            cell.margin_top=Inches(0.02); cell.margin_bottom=Inches(0.02)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size=Pt(font); r.font.name="Calibri"

def set_cell(cell, text, color=TEXT, bold=False, align=PP_ALIGN.LEFT, size=11):
    tf=cell.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    for r in list(p.runs): r.text=""
    run(p, text, size, color, bold)

# =====================================================================
# SLIDE 1 — INTRO
# =====================================================================
s = slide()
line1(s, 0.9, 1.35, 11, "MILESTONE REPORT", 15, ORANGE, bold=True)
box, tf = tb(s, 0.85, 1.7, 11.6, 1.3)
p=tf.paragraphs[0]; run(p, "SnapChef ", 64, TEXT, bold=True); run(p, "🥕", 54, TEAL, bold=True)
box, tf = tb(s, 0.9, 3.0, 9.6, 1.1)
p=tf.paragraphs[0]
run(p, "A stick-on smart-fridge companion that recognizes groceries, tracks "
       "your inventory, and suggests recipes — on two ESP32-S3 boards.", 21, SUB)

c1 = card(s, 0.9, 4.5, 3.7, 1.9)
_,tf=tb(s,1.1,4.7,3.3,1.6)
run(tf.paragraphs[0], "Team", 19, TEXT, bold=True)
p=tf.add_paragraph(); p.space_before=Pt(8); run(p,"Team #  ",15,SUB); ph(p,"FILL IN")

c2 = card(s, 4.8, 4.5, 3.7, 1.9)
_,tf=tb(s,5.0,4.7,3.4,1.7)
run(tf.paragraphs[0], "Members", 19, TEXT, bold=True)
p=tf.add_paragraph(); p.space_before=Pt(6); run(p,"Menglh ",14,TEXT); ph(p,"full name")
p=tf.add_paragraph(); run(p,"Iris (Yzh) ",14,TEXT); ph(p,"full name")
p=tf.add_paragraph(); ph(p,"+ other members?")

c3 = card(s, 8.7, 4.5, 3.7, 1.9)
_,tf=tb(s,8.9,4.7,3.4,1.6)
run(tf.paragraphs[0], "Repository", 19, TEXT, bold=True)
p=tf.add_paragraph(); p.space_before=Pt(8); run(p,"github.com/IrisYzh/SnapChef",14,TEAL,bold=True)
footer(s,1)

# =====================================================================
# SLIDE 2 — PROBLEM
# =====================================================================
s = slide()
kicker_title(s, "Problem", "Food gets wasted because no one tracks the fridge")
box, tf = tb(s, 0.6, 1.7, 6.6, 4.6)
p=tf.paragraphs[0]
run(p,"People forget what they already own, buy duplicates, and let produce rot "
      "out of sight. Existing “smart fridges” are expensive, built-in, and can’t be "
      "added to the fridge you already have.", 17, TEXT)
bx = bullets(s, 0.6, 3.3, 6.6, 3.0, [
 ("No lightweight way to log items in/out without an app or typing.", False),
 ("Produce is hard to track — no barcode, many varieties.", False),
 ("Receipts hold the data but get thrown away.", False),
 ("Users want recipe ideas from what they own before it spoils.", False),
], size=16)
c = card(s, 7.5, 1.7, 5.2, 4.4)
_,tf = tb(s, 7.8, 1.95, 4.7, 4.0)
run(tf.paragraphs[0], "Intended users", 19, TEAL, bold=True)
for txt in ["Home cooks & busy households who shop weekly and cook from fresh produce.",
            "Students / shared apartments who can’t justify a built-in smart fridge.",
            "People who want zero-friction input — hold an item to a camera or snap a receipt."]:
    p=tf.add_paragraph(); p.space_before=Pt(10)
    run(p,"•  ",16,ORANGE,bold=True); run(p,txt,16,TEXT)
p=tf.add_paragraph(); p.space_before=Pt(14)
run(p,"A peel-and-stick device that turns any fridge into a tracked one.",13,SUB,italic=True)
footer(s,2)

# =====================================================================
# SLIDE 3 — SOLUTION
# =====================================================================
s = slide()
kicker_title(s, "Solution", "A two-board fridge sticker: camera + touchscreen")
line1(s, 0.6, 1.55, 12.1,
      "A sensor/compute board (camera, distance sensor, Wi-Fi) talks over ESP-NOW to a "
      "4.3\" touch display that owns the UI and inventory. Heavy AI runs in the cloud.",
      14.5, SUB, h=0.6)
flows = [("🥦  Put in","Hold a fruit/veg to the camera, or snap a grocery receipt. Items are recognized and merged into your fridge inventory."),
         ("🍳  Take out + recipes","Scan an item to remove it, or pick from the list. SnapChef suggests recipes from what's left and shows which items each uses."),
         ("📋  My Fridge","Browse inventory sorted by name & time added; remove by tap. Proximity wake lights the screen as you approach.")]
for i,(t,d) in enumerate(flows):
    x = 0.6 + i*4.07
    c = card(s, x, 2.25, 3.8, 1.85)
    _,tf=tb(s,x+0.22,2.45,3.4,1.5)
    run(tf.paragraphs[0],t,17,TEXT,bold=True)
    p=tf.add_paragraph(); p.space_before=Pt(6); run(p,d,12.5,SUB)
feat = card(s, 0.6, 4.35, 6.0, 2.35)
_,tf=tb(s,0.85,4.55,5.5,2.0)
run(tf.paragraphs[0],"Key features",17,TEAL,bold=True)
for txt in ["On-device ML (TFLM, 10 produce classes) — instant, offline.",
            "Cloud LLM-vision fallback for any produce outside those classes.",
            "Receipt OCR → structured grocery list with refrigeration tags.",
            "LLM recipe suggestions from current inventory."]:
    p=tf.add_paragraph(); p.space_before=Pt(5); run(p,"•  ",13,ORANGE,bold=True); run(p,txt,13,TEXT)
hw = card(s, 6.8, 4.35, 5.9, 2.35)
_,tf=tb(s,7.05,4.55,5.4,2.0)
run(tf.paragraphs[0],"Hardware & stack",17,TEAL,bold=True)
for txt in ["Main: Seeed XIAO ESP32S3 Sense (camera + HC-SR04 + Wi-Fi).",
            "Display: Waveshare ESP32-S3-Touch-LCD-4.3 (LVGL UI).",
            "Link: ESP-NOW peer-to-peer between boards.",
            "Backend: FastAPI on Railway + Claude Vision / Baidu / Textract."]:
    p=tf.add_paragraph(); p.space_before=Pt(5); run(p,"•  ",13,ORANGE,bold=True); run(p,txt,13,TEXT)
footer(s,3)

# =====================================================================
# SLIDE 4 — ARCHITECTURE
# =====================================================================
s = slide()
kicker_title(s, "Updated Architecture", "System diagram — color = maturity · line = interface")
# legend
lx=0.6; ly=1.45
def swatch(x, col, label):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x),Inches(ly+0.04),Inches(0.18),Inches(0.18))
    sh.fill.solid(); sh.fill.fore_color.rgb=col; sh.line.fill.background(); sh.shadow.inherit=False
    line1(s, x+0.26, ly, 1.7, label, 11.5, SUB)
swatch(0.6, GREEN_B, "Proven")
swatch(2.0, YEL_B, "In progress")
swatch(3.7, RED_B, "Blocked")
line1(s, 5.4, ly, 3.0, "——  tested interface", 11.5, SUB)
line1(s, 8.2, ly, 3.5, "– – –  unverified interface", 11.5, SUB)

# column headers — DISPLAY · MAIN · CLOUD (main in the middle = the data hub,
# so every link is between adjacent columns and stays in the empty gaps)
line1(s, 0.6, 1.95, 3.7, "DISPLAY · Waveshare LVGL", 12.5, ORANGE, bold=True, align=PP_ALIGN.CENTER)
line1(s, 4.8, 1.95, 3.7, "MAIN · XIAO ESP32S3 Sense", 12.5, ORANGE, bold=True, align=PP_ALIGN.CENTER)
line1(s, 9.0, 1.95, 3.7, "CLOUD · FastAPI @ Railway", 12.5, ORANGE, bold=True, align=PP_ALIGN.CENTER)

MW=3.7; CW=3.7
dx=0.6; mx=4.8; cx=9.0
# DISPLAY (left)
block(s, dx,2.35,CW,0.74,"LVGL UI + screen flows","put-in / take-out / fridge","g")
block(s, dx,3.20,CW,0.74,"Inventory v2 (NVS)","timestamps · normalized names","g")
block(s, dx,4.05,CW,0.74,"Manual take-out + recipe UI","used/missing remove chips","y")
block(s, dx,4.90,CW,0.64,"Display E2E verified","2 consecutive runs ✓","g")
# MAIN (center)
block(s, mx,2.35,MW,0.74,"Camera capture (JPEG)","UXGA · FB-OVF-safe","g")
block(s, mx,3.20,MW,0.74,"TFLM classifier (10 classes)","on-device · conf≥0.90","g")
block(s, mx,4.05,MW,0.74,"Cloud-fallback controller","5 s local → LLM call","g")
block(s, mx,4.90,MW,0.64,"HC-SR04 proximity wake","","g")
block(s, mx,5.66,MW,0.64,"Wi-Fi + NTP clock","","g")
# CLOUD (right)
block(s, cx,2.35,CW,0.70,"/produce/recognize-llm","Claude Vision · 200 OK","g")
block(s, cx,3.15,CW,0.70,"/receipts/analyze","Textract + Claude","g")
block(s, cx,3.95,CW,0.70,"/recipes/list · /steps","used_fridge / missing","y")
block(s, cx,4.75,CW,0.62,"/healthz","","g")

# connectors — all between adjacent columns, routed through the empty gaps
GAP_L = 4.55   # centre of display|main gap
GAP_R = 8.75   # centre of main|cloud gap
# ESP-NOW display<->main (solid, proven)
connector(s, dx+CW, 2.62, mx, 2.62, color=GREEN, width=2.25)        # display->main
connector(s, mx, 2.92, dx+CW, 2.92, color=GREEN, width=2.25)        # main->display
line1(s, GAP_L-0.65, 2.16, 1.3, "ESP-NOW", 10, SUB, align=PP_ALIGN.CENTER)
# time_sync main->display (solid, proven — display E2E now works)
connector(s, mx, 5.95, dx+CW, 3.57, color=GREY, width=1.8)
line1(s, GAP_L-0.65, 5.52, 1.3, "time_sync", 9.5, SUB, align=PP_ALIGN.CENTER)
# main->cloud HTTPS (produce / receipts) — solid, tested
connector(s, mx+MW, 2.95, cx, 2.90, color=GREY, width=1.9)
line1(s, GAP_R-0.85, 2.50, 1.7, "HTTPS", 9.5, SUB, align=PP_ALIGN.CENTER)
# main->cloud recipes — dashed, unverified
connector(s, mx+MW, 4.10, cx, 4.05, color=GREY, width=1.9, dash=True)
line1(s, GAP_R-0.85, 4.20, 1.7, "recipes", 9.5, SUB, align=PP_ALIGN.CENTER)

line1(s, 0.6, 6.62, 12.1,
  "Change since Milestone 1: BLE link → ESP-NOW; on-device-only recognition now has a cloud LLM fallback; "
  "inventory gained timestamps, normalized names, manual take-out and recipe used/missing UI.",
  11.5, RGBColor(0x77,0x77,0x77))
footer(s,4)

# =====================================================================
# SLIDE 5 — OPTIMIZATION
# =====================================================================
s = slide()
kicker_title(s, "Optimization · since Milestone 2", "Two measurable improvements")

def opt_table(x, title, rows):
    line1(s, x+0.05, 1.55, 6.0, title, 17, TEXT, bold=True)
    t = s.shapes.add_table(len(rows), 3, Inches(x), Inches(1.95), Inches(5.9), Inches(2.2)).table
    t.columns[0].width=Inches(2.8); t.columns[1].width=Inches(1.55); t.columns[2].width=Inches(1.55)
    for ri,(a,b,c,cls) in enumerate(rows):
        set_cell(t.cell(ri,0), a, ORANGE if ri==0 else TEXT, bold=(ri==0), size=11.5)
        col = TEXT if ri==0 else (GREEN if cls=="g" else RED if cls=="r" else TEXT)
        set_cell(t.cell(ri,1), b, ORANGE if ri==0 else (RED if cls in("r","b") else SUB), bold=(ri==0), align=PP_ALIGN.CENTER, size=11.5)
        set_cell(t.cell(ri,2), c, ORANGE if ri==0 else GREEN, bold=(ri==0 or ri>0), align=PP_ALIGN.CENTER, size=11.5)
    style_table(t, font=11.5)
    return t

opt_table(0.6, "①  Recognition coverage & accuracy", [
  ("Metric","Before (M2)","After","h"),
  ("Recognizable items","10 classes","arbitrary","g"),
  ("Plastic / hand-held / glare","Uncertain 0.37","correct 0.95","g"),
  ("Lock threshold","0.80","0.90","g"),
  ("Out-of-class result","\"Unknown\"","named (LLM)","g"),
])
opt_table(6.8, "②  Latency & capture quality", [
  ("Metric","Before (M2)","After","h"),
  ("Local window → fallback","20 s","5 s","g"),
  ("Cloud recognize latency","—","1.2–2.7 s","g"),
  ("Capture quality (lower=better)","q12","q6","g"),
  ("Camera overflow (FB-OVF)","on q<init","eliminated","g"),
])
c = card(s, 0.6, 4.55, 12.1, 2.15)
_,tf=tb(s,0.85,4.72,11.6,1.9)
run(tf.paragraphs[0],"Qualitative analysis",16,TEAL,bold=True)
for txt in [
 ("Coverage: ","the on-device model knew only 10 classes and rejected the rest; the LLM fallback makes recognition open-vocabulary while keeping the fast offline path for common items."),
 ("Robustness: ","on hard inputs Baidu scored 0.37→Uncertain while Claude Vision returned the correct item at 0.95; the 0.90 lock also cuts false on-device locks."),
 ("Responsiveness: ","cutting the local window 20s→5s reaches a cloud answer sooner for off-list items, and higher-quality capture improves OCR + produce accuracy with no overflow."),
]:
    p=tf.add_paragraph(); p.space_before=Pt(5)
    run(p,"•  ",12.5,ORANGE,bold=True); run(p,txt[0],12.5,TEXT,bold=True); run(p,txt[1],12.5,SUB)
footer(s,5)

# =====================================================================
# SLIDE 6 — READINESS
# =====================================================================
s = slide()
kicker_title(s, "Readiness", "End-to-end runs & cross-condition testing")
# E2E table
line1(s, 0.6, 1.6, 6.2, "E2E pipeline — 2 consecutive runs (full device)", 16, TEXT, bold=True)
line1(s, 0.6, 1.95, 6.4, "camera → ESP-NOW → cloud → result shown on the display", 12, SUB)
t = s.shapes.add_table(3,4, Inches(0.6),Inches(2.35),Inches(5.9),Inches(1.4)).table
t.columns[0].width=Inches(0.8);t.columns[1].width=Inches(2.6);t.columns[2].width=Inches(1.1);t.columns[3].width=Inches(1.4)
e2e=[("Run","Input","Latency","Result"),
     ("1","OOD corn (off-list)","2.1 s","Corn 0.95 ✓"),
     ("2","bell pepper","1.3 s","Pepper 0.85 ✓")]
for ri,(a,b,c,d) in enumerate(e2e):
    set_cell(t.cell(ri,0),a,ORANGE if ri==0 else TEXT,bold=(ri==0),size=11.5)
    set_cell(t.cell(ri,1),b,ORANGE if ri==0 else TEXT,bold=(ri==0),size=11.5)
    set_cell(t.cell(ri,2),c,ORANGE if ri==0 else SUB,bold=(ri==0),align=PP_ALIGN.RIGHT,size=11.5)
    set_cell(t.cell(ri,3),d,ORANGE if ri==0 else GREEN,bold=(ri==0),align=PP_ALIGN.CENTER,size=11.5)
style_table(t,font=11.5)
line1(s, 0.6, 3.95, 6.3, "Display device now built & running — both runs completed the full pipeline and rendered on-screen.", 11, SUB)

# cross condition table
line1(s, 6.9, 1.6, 6.0, "Cross-condition results", 16, TEXT, bold=True)
t2 = s.shapes.add_table(6,3, Inches(6.9),Inches(1.95),Inches(5.8),Inches(2.7)).table
t2.columns[0].width=Inches(2.0);t2.columns[1].width=Inches(2.0);t2.columns[2].width=Inches(1.8)
cc=[("Condition","Input","Outcome","h"),
    ("Real produce","ripe tomato","Tomato 0.88 ✓","g"),
    ("OOD produce","corn (off-list)","Corn 0.95 ✓","g"),
    ("Hand-held / glare","bell pepper","Pepper 0.85 ✓","g"),
    ("Non-produce","water bottle","Okra 0.75 ✗","r"),
    ("Diff device / user","2 boards, 2 IPs","both 200 ✓","g")]
for ri,(a,b,c,cls) in enumerate(cc):
    set_cell(t2.cell(ri,0),a,ORANGE if ri==0 else TEXT,bold=(ri==0),size=11)
    set_cell(t2.cell(ri,1),b,ORANGE if ri==0 else SUB,bold=(ri==0),size=11)
    col = TEXT if ri==0 else (GREEN if cls=="g" else (RED if cls=="r" else YELLOW))
    set_cell(t2.cell(ri,2),c,ORANGE if ri==0 else col,bold=(ri==0),size=11)
style_table(t2,font=11)
line1(s, 6.9, 4.78, 5.8, "OOD = outside the 10 on-device classes — handled by the LLM fallback. Names normalized: Red Pepper → Pepper.", 10, SUB)

c = card(s, 0.6, 5.35, 12.1, 1.35)
_,tf=tb(s,0.85,5.5,11.6,1.1)
run(tf.paragraphs[0],"Status summary",15,TEAL,bold=True)
p=tf.add_paragraph(); p.space_before=Pt(4)
run(p,"Proven: ",12.5,GREEN,bold=True); run(p,"full device E2E (camera→cloud→display), ESP-NOW link, produce-LLM & receipt endpoints, on-device TFLM, inventory v2.   ",12.5,TEXT)
p=tf.add_paragraph(); p.space_before=Pt(2)
run(p,"In progress: ",12.5,YELLOW,bold=True); run(p,"recipe used/missing UI & /recipes endpoints.   ",12.5,TEXT)
run(p,"Known gap: ",12.5,ORANGE,bold=True); run(p,"non-produce over-recognition → backend Uncertain rules + vocab.",12.5,TEXT)
footer(s,6)

# =====================================================================
# SLIDE 7 — BUDGET
# =====================================================================
s = slide()
kicker_title(s, "Budget Update", "Spend to date")
box,tf = tb(s,0.6,1.5,12,0.4)
p=tf.paragraphs[0]; run(p,"Unit prices are typical retail estimates — ",12.5,SUB); ph(p,"replace with actual receipts")

rows=[("Item","Qty","Unit (est.)","Subtotal","Status","h"),
 ("Seeed XIAO ESP32S3 Sense (incl. camera)","1","$23.99","$23.99","purchased","g"),
 ("Waveshare ESP32-S3-Touch-LCD-4.3","1","$45.99","$45.99","purchased","g"),
 ("HC-SR04 ultrasonic sensor","1","$3.50","$3.50","purchased","g"),
 ("Wires, resistor divider, misc.","1","$5.00","$5.00","purchased","g"),
 ("Railway hosting (backend)","~1 mo","$5.00","$5.00","recurring","y"),
 ("Claude / Baidu / Textract API","usage","~$0–10","~$5.00","usage","y"),
 ("Total spent (est.)","","","≈ $88.48","","t")]
t = s.shapes.add_table(len(rows),5, Inches(0.6),Inches(2.05),Inches(12.1),Inches(3.1)).table
for w,c in zip([5.6,1.0,1.7,1.7,2.1],range(5)): t.columns[c].width=Inches(w)
for ri,r in enumerate(rows):
    a,b,c,d,e,cls=r
    set_cell(t.cell(ri,0),a,ORANGE if ri==0 else TEXT,bold=(ri==0 or cls=="t"),size=12)
    set_cell(t.cell(ri,1),b,ORANGE if ri==0 else SUB,bold=(ri==0),align=PP_ALIGN.RIGHT,size=12)
    set_cell(t.cell(ri,2),c,ORANGE if ri==0 else SUB,bold=(ri==0),align=PP_ALIGN.RIGHT,size=12)
    set_cell(t.cell(ri,3),d,ORANGE if ri==0 else (TEAL if cls=="t" else TEXT),bold=(ri==0 or cls=="t"),align=PP_ALIGN.RIGHT,size=12)
    stcol = GREEN if cls=="g" else (YELLOW if cls=="y" else TEXT)
    set_cell(t.cell(ri,4),e,ORANGE if ri==0 else stcol,bold=(ri==0),size=12)
style_table(t,font=12)

for i,(title,val,col) in enumerate([("Total budget","$___",None),("Spent","≈ $88",ORANGE),("Remaining","$___",None)]):
    x=0.6+i*4.07
    card(s,x,5.45,3.8,1.2)
    line1(s,x+0.25,5.6,3.3,title,15,TEXT,bold=True)
    box,tf=tb(s,x+0.25,5.95,3.3,0.7)
    p=tf.paragraphs[0]
    if col: run(p,val,30,col,bold=True)
    else: ph(p,val)
line1(s,0.6,6.78,12,"Fill in your allocated total budget; remaining = total − spent.",11,RGBColor(0x77,0x77,0x77))
footer(s,7)

out = "/Users/menglh/Desktop/GIX26Spring/SnapChef/docs/SnapChef_Milestone_Report.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
